"""과일 카운팅 프로젝트(Plans 1-4)의 이론적 배경 + 핵심 인사이트를
정리한 연구용 PowerPoint 생성.

초점: 이론(SAM · CLIP · PseCo 구조) + 인사이트(파인튜닝·좌표계·CLIP 정렬
취약성 등 실험적 발견). 수치/결론이 바뀌면 스크립트 수정 후 재실행으로
즉시 반영.

실행: python docs/presentations/generate_insights_ppt.py
출력: docs/presentations/fruit_counting_insights.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

OUT = Path(__file__).resolve().parent / "fruit_counting_insights.pptx"

# 차분한 연구용 색상 팔레트
C_TITLE = RGBColor(0x1F, 0x3B, 0x5C)       # 네이비
C_HEADER = RGBColor(0x2E, 0x5F, 0x8A)      # 스틸 블루
C_BODY = RGBColor(0x2A, 0x2A, 0x2A)        # 거의 블랙
C_ACCENT = RGBColor(0xC2, 0x56, 0x2D)      # 러스트
C_MUTED = RGBColor(0x6B, 0x6B, 0x6B)       # 그레이


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(12), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = C_TITLE

    tb2 = slide.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(12), Inches(1.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.size = Pt(20)
    r2.font.color.rgb = C_MUTED

    tb3 = slide.shapes.add_textbox(Inches(0.7), Inches(6.8), Inches(12), Inches(0.5))
    tf3 = tb3.text_frame
    r3 = tf3.paragraphs[0].add_run()
    r3.text = "내부 연구 노트 · 2026-04-22"
    r3.font.size = Pt(11)
    r3.font.color.rgb = C_MUTED
    r3.font.italic = True


def add_section_slide(prs, section_number, section_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(12), Inches(1.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"Part {section_number}"
    r.font.size = Pt(16)
    r.font.color.rgb = C_MUTED

    p = tf.add_paragraph()
    r = p.add_run()
    r.text = section_title
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = C_TITLE


def add_content_slide(prs, title, bullets, footer=None):
    """bullets: [(level:int, text:str), ...] 또는 단순 [str, ...]."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.8))
    tf = tb.text_frame
    r = tf.paragraphs[0].add_run()
    r.text = title
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = C_HEADER

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(12.0), Inches(5.8))
    bt = body.text_frame
    bt.word_wrap = True
    first = True
    for item in bullets:
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        if first:
            p = bt.paragraphs[0]
            first = False
        else:
            p = bt.add_paragraph()
        p.level = level
        r = p.add_run()
        r.text = ("• " if level == 0 else "– ") + text
        r.font.size = Pt(18 if level == 0 else 15)
        r.font.color.rgb = C_BODY if level == 0 else C_MUTED

    if footer:
        tb2 = slide.shapes.add_textbox(Inches(0.7), Inches(7.0), Inches(12), Inches(0.4))
        tf2 = tb2.text_frame
        r2 = tf2.paragraphs[0].add_run()
        r2.text = footer
        r2.font.size = Pt(11)
        r2.font.italic = True
        r2.font.color.rgb = C_MUTED


def add_table_slide(prs, title, headers, rows, footer=None, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.8))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = title
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = C_HEADER

    n_rows = len(rows) + 1
    n_cols = len(headers)
    left, top = Inches(0.7), Inches(1.3)
    width = Inches(12.0)
    height = Inches(min(5.5, 0.5 + 0.45 * n_rows))

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_HEADER

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(12)
            run.font.color.rgb = C_BODY
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)

    if footer:
        tb2 = slide.shapes.add_textbox(Inches(0.7), Inches(7.0), Inches(12), Inches(0.4))
        r2 = tb2.text_frame.paragraphs[0].add_run()
        r2.text = footer
        r2.font.size = Pt(11)
        r2.font.italic = True
        r2.font.color.rgb = C_MUTED


def add_two_col_slide(prs, title, left_heading, left_bullets,
                      right_heading, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.8))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = title
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = C_HEADER

    def _col(left_in, heading, items):
        tb = slide.shapes.add_textbox(
            Inches(left_in), Inches(1.3), Inches(6.1), Inches(5.8)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        hr = p.add_run()
        hr.text = heading
        hr.font.size = Pt(18)
        hr.font.bold = True
        hr.font.color.rgb = C_ACCENT
        for it in items:
            if isinstance(it, tuple):
                level, text = it
            else:
                level, text = 0, it
            p = tf.add_paragraph()
            p.level = level
            r2 = p.add_run()
            r2.text = ("• " if level == 0 else "– ") + text
            r2.font.size = Pt(14 if level == 0 else 12)
            r2.font.color.rgb = C_BODY if level == 0 else C_MUTED

    _col(0.7, left_heading, left_bullets)
    _col(7.0, right_heading, right_bullets)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # === 표지 ===
    add_title_slide(
        prs,
        "PseCo 기반 과일 카운팅:\nCLIP 정렬 분류기에 대한 사례 연구",
        "네 번의 학습·진단·프로덕션 이터레이션에서 얻은 "
        "이론적 배경과 핵심 인사이트 — 과수원 이미지 기반",
    )

    # =============================================================
    # Part 1: 문제 · 아키텍처 · 이론
    # =============================================================
    add_section_slide(prs, 1, "문제 · 아키텍처 · 이론")

    add_content_slide(
        prs,
        "문제 정의",
        [
            "목표: 배포 카메라로 촬영한 과수원 이미지에서 보호 봉지 씌운 과일의 개수를 셈.",
            "원본 파이프라인: 5단계 구성 (Deblur → PseCo → Super-Resolution → Classifier → Defect).",
            (1, "각 단계는 외부 모델의 얇은 래퍼였고, 추론 전용 — 학습 코드 없음."),
            (1, "자체 과수원 라벨 없음, 체계적 평가 없음."),
            "우리 과제: 전면 재설계 + 학습 코드 추가 + 실제 과수원 이미지로 검증.",
            "연구 질문: 오픈월드 카운팅 모델을 일반성을 잃지 않고 도메인에 적응시킬 수 있는가?",
        ],
        footer="스펙: docs/superpowers/specs/2026-04-19-fruit-counting-redesign-design.md",
    )

    add_content_slide(
        prs,
        "이론적 구성 요소 (1/3) — SAM (Segment Anything Model)",
        [
            "Meta AI, 2023 — 클래스 무관 세그멘테이션 파운데이션 모델.",
            (1, "ViT-H 이미지 인코더: 1024×1024 입력 → (256, 64, 64) 피처 맵."),
            (1, "프롬프트 기반 마스크 디코더: 점 / 박스 / 텍스트 프롬프트 수용."),
            (1, "10억 개 이상 마스크(SA-1B)로 학습 — 미지 도메인에서도 강한 일반화."),
            "카운팅에서의 역할: 후보 점이 주어지면 SAM이 해당 객체의 제안(proposal)을 생성.",
            (1, "`forward_sam_with_embeddings(features, points=...)` — 점당 k=5개의 가능한 마스크 반환."),
            (1, "anchor box(점 ± half_side)로 마스크 가설을 정제."),
            "핵심 성질: SAM은 \"그것이 무엇인지\"는 모르고, \"객체스러운 덩어리\"의 위치만 안다.",
        ],
    )

    add_content_slide(
        prs,
        "이론적 구성 요소 (2/3) — CLIP 텍스트 정렬",
        [
            "OpenAI, 2021 — 4억 (이미지, 캡션) 페어로 이미지·텍스트 인코더를 공동 학습.",
            (1, "텍스트 인코더: 클래스 이름 → 512-d 임베딩 (ViT-B/32 백본 기준)."),
            (1, "이미지 인코더: 크롭을 동일한 512-d 공간에 매핑."),
            (1, "공유 공간에서의 코사인 유사도로 zero-shot 분류."),
            "PseCo에서의 역할: ROIHeadMLP의 512-d 출력을 CLIP 텍스트 피처와 비교해 분류.",
            (1, "dot(ROI_feature, text_feature) → 점수. 높을수록 클래스 이름과 매칭."),
            "미묘한 지점: OpenAI CLIP 가중치는 QuickGELU 활성화로 학습됨 (표준 GELU 아님).",
            (1, "`open_clip.create_model('ViT-B-32')` 기본값은 표준 GELU → 텍스트 피처가 조용히 drift."),
            (1, "반드시 `'ViT-B-32-quickgelu'` 변형을 써야 함 — Plan 3에서 이것을 놓쳐 큰 대가를 치름."),
        ],
    )

    add_content_slide(
        prs,
        "이론적 구성 요소 (3/3) — PseCo 파이프라인 (ICLR 2024)",
        [
            "\"Point, Segment, and Count\" (Huang et al., 2024).",
            "캐시된 SAM 피처 위에서 3단계 추론:",
            (1, "PointDecoder: dense heatmap → top-K 후보 점 (자체 NMS 적용)."),
            (1, "SAM을 proposal 생성기로 사용: 각 점에 대해 5개의 마스크·박스 가설 제안."),
            (1, "ROIHeadMLP + CLIP: 각 (영역, 텍스트) 쌍을 점수화 → NMS → threshold."),
            "카운팅에 효과적인 이유:",
            (1, "클래스 무관한 SAM이 과수원 같은 새 도메인을 인코더 재학습 없이 커버."),
            (1, "텍스트 프롬프트가 \"셀 대상\"을 지정 — \"apples\" → \"bread rolls\" 전환 시 재학습 불필요."),
            (1, "분리 구조: 위치 추정(SAM+PointDecoder)은 공유, 분류(MLP+CLIP)는 교체 가능."),
            "Upstream 학습 목표: pseudo-box IoU에 대한 BCE + CLIP 일관성 손실(정렬 유지용).",
        ],
        footer="github.com/Hzzone/PseCo — FSC-147 test split에서 MAE 16–20 보고",
    )

    # =============================================================
    # Part 2: 우리의 실험 이터레이션
    # =============================================================
    add_section_slide(prs, 2, "FSC-147 + 과수원 위에서의 네 번의 이터레이션")

    add_content_slide(
        prs,
        "우리의 Plan 타임라인",
        [
            "Plan 1 (Foundation): 스캐폴드, CLI, 설정, 데이터 진단. 학습 없음.",
            "Plan 2 (학습 인프라): SAM 피처 캐시 (FSC-147 4945장의 fp16 임베딩 8.7 GB), Runner, TensorBoard, 체크포인트. Placeholder objective.",
            "Plan 3 (첫 실제 파인튜닝): CLIP 텍스트 프롬프트 + 점-박스 포함 기반 pos/neg 레이블.",
            (1, "발견한 버그 3종: CUDA device mismatch, 좌표계 mismatch, QuickGELU 변형 mismatch."),
            (1, "K=32 proposal로 FSC-147 val/mae = 47.62 달성. \"predict-zero\" baseline 50.58 돌파."),
            (1, "이 시점에서는 성공처럼 보였음."),
            "Plan 4 (과수원 실전 검증): 직접 센 과수원 이미지 4장이 Plan 3의 숨은 실패를 드러냄.",
            (1, "파인튜닝 포기 → upstream MLP + 프롬프트 엔지니어링 + SAM anchor 튜닝으로 피봇."),
            (1, "최종 과수원 MAE 11.75 — 파인튜닝 없이 달성."),
        ],
    )

    add_table_slide(
        prs,
        "이터레이션별 수치",
        headers=[
            "버전", "FSC-147 val/mae",
            "과수원 MAE (4장)", "교훈",
        ],
        rows=[
            ("Plan 2 (placeholder)", "50.58 (축퇴)", "—", "항상 0만 예측"),
            ("Plan 3 v1 (no fix)", "62.90", "—", "좌표계 버그"),
            ("Plan 3 v2 (+ QuickGELU)", "62.79", "—", "미미"),
            ("Plan 3 v3 (+ 좌표 수정)", "54.26", "—", "좌표 수정이 지배적"),
            ("Plan 3 v4 (K=32)", "47.62", "400+ ❌", "프로덕션에서 prompt-agnostic"),
            ("Plan 4 upstream MLP (기본)", "—", "14.25", "학습 0, 충분히 경쟁적"),
            ("Plan 4 + threshold 튜닝", "—", "13.50", "그리드 서치 소폭 개선"),
            ("Plan 4 + anchor_size=24", "—", "11.75 ⭐", "스케일이 핵심"),
        ],
        col_widths=[3.4, 2.4, 2.7, 3.4],
    )

    # =============================================================
    # Part 3: 실패 모드에서 얻은 인사이트
    # =============================================================
    add_section_slide(prs, 3, "실패 모드에서 얻은 인사이트")

    add_content_slide(
        prs,
        "인사이트 1 — \"Predict Zero\"는 기만적 baseline",
        [
            "카운팅 데이터셋에서 모든 이미지에 0 예측 시 MAE = E[GT count].",
            (1, "FSC-147 val 기준 평균 GT 개수 ≈ 50.58."),
            (1, "Plan 2의 placeholder 트레이너가 매 에폭 `val/mae = 50.58`을 유지."),
            (1, "순진한 해석: \"자명하지 않은 손실에서 정체\". 실제로는 모델이 붕괴한 상태."),
            "Plan 3의 `val/mae = 47.62` — zero-prediction 대비 겨우 한 자릿수 개선.",
            (1, "처음에는 baseline 돌파로 축하했지만, 사실상 여전히 zero-prediction 오차의 93% 수준."),
            "연구적 함의: 카운팅에서는 predict-zero와 predict-mean 두 기준선을 모두 리포트해야 함.",
            (1, "predict-zero를 이기지만 predict-mean을 못 이기면, 단순히 데이터셋 평균을 학습한 것일 수 있음."),
            (1, "수렴 지표는 이 자명한 기준선을 명시적으로 차감한 형태여야 함."),
        ],
    )

    add_content_slide(
        prs,
        "인사이트 2 — CLIP 정렬은 순진한 파인튜닝에 취약함",
        [
            "Plan 3 학습: 이미지당 하나의 클래스, proposal의 (pos/neg) 레이블에 대한 binary cross-entropy.",
            "결과: 파인튜닝된 MLP가 프롬프트에 무감각해짐.",
            (1, "같은 사과 이미지에서 `apples`, `apple in paper bag`, `protective fruit bag` 모두 ~500개 예측."),
            (1, "Upstream MLP는 동일 이미지에서 6 / 25 / 8로 프롬프트를 구분."),
            "메커니즘: 단일 클래스 binary CE 기울기는 \"시각 특징이 pos/neg\"만 학습, 텍스트 정렬 신호는 무시됨.",
            (1, "그래디언트가 'CLIP 텍스트 일치/불일치'를 한 번도 처벌하지 않음."),
            (1, "Pretrained의 CLIP 정렬이 파인튜닝 동안 서서히 decay."),
            "Upstream PseCo는 보조 CLIP 일관성 손실(`cls_loss2`)을 씀 — 우리는 생략했고 결과는 위와 같음.",
            (1, "연구 질문: 파인튜닝 동안 CLIP 정렬 손실을 정량화해 조기 중단 기준으로 쓸 수 있는가?"),
            (1, "후보 지표: 고정된 triplet (이미지, 정답 클래스, 오답 클래스)에 대한 prompt-discrimination 점수."),
        ],
    )

    add_content_slide(
        prs,
        "인사이트 3 — 좌표계 버그가 학습을 조용히 낭비시킴",
        [
            "Plan 3의 레이블은 FSC-147 점 어노테이션(원본 이미지 좌표계, 긴 변 384 px) 기반.",
            "PointDecoder 예측은 SAM의 1024-px padding 좌표계.",
            (1, "점-박스 포함 체크가 두 좌표계를 같은 프레임으로 가정해 비교."),
            (1, "포함 판정이 거의 일어나지 않음 → 대부분 proposal이 negative로 분류."),
            "증상: 학습 loss는 내려감 (모두 negative라고 답하면 쉬움), val/mae는 62-63에서 정체.",
            (1, "문제를 \"쉽게 학습 가능하게 만들어서 loss를 줄이는\" 버그."),
            "수정: GT 점을 `1024 / max(H, W)` 배로 1024 공간에 스케일링. 한 커밋으로 val/mae 62.79 → 54.26.",
            "연구·엔지니어링 교훈:",
            (1, "멀티모듈 파이프라인은 좌표 프레임을 조용히 섞음. 프레임 일치를 주장하는 단일 테스트를 써야 함."),
            (1, "\"loss가 내려감\" 자체는 절대로 정확성의 증거가 아님."),
        ],
    )

    add_content_slide(
        prs,
        "인사이트 4 — QuickGELU vs GELU: 활성화 0.3% 차이, 의미 100% drift",
        [
            "OpenAI의 원본 CLIP은 QuickGELU로 학습됨: `x * sigmoid(1.702 * x)`.",
            "`open_clip.create_model('ViT-B-32', pretrained='openai')` 기본값은 표준 GELU.",
            (1, "텍스트 임베딩은 L2 기준 ~2%만 바뀜."),
            (1, "하지만 QuickGELU 임베딩 위에서 학습된 MLP 헤드와 스펙이 맞지 않게 됨."),
            "Upstream PseCo `MLP_small_box_w1_zeroshot.tar` 은 QuickGELU CLIP 피처로 학습됨.",
            (1, "우리는 경고를 놓친 채 표준-GELU 피처로 2번의 학습 런을 돌림."),
            (1, "이 수정만으로 CLIP 정렬을 회복하진 못했지만(인사이트 2 참조), 필요조건임."),
            "연구·엔지니어링 교훈:",
            (1, "사전학습 모델의 activation variant 태그는 semver-breaking으로 취급해야 함. 체크포인트 메타에 명시."),
        ],
    )

    add_content_slide(
        prs,
        "인사이트 5 — SAM anchor size는 객체 스케일에 맞아야 함",
        [
            "PseCo 추론은 anchor box(점 ± half_side)를 SAM의 box-refinement 헤드 프롬프트로 사용.",
            "Upstream demo: 1024-padding 입력에서 anchor_size = 8 (약 8 px half-side, ~16 px full side).",
            (1, "FSC-147에는 적합 (이미지 긴 변 384 px, 객체 ~10–30 px)."),
            (1, "폰으로 찍은 과수원 사진(1024 공간에서 객체 ~80–120 px)에서는 작음."),
            "anchor_size ∈ {4, 8, 12, 16, 24, 32} 그리드 서치: 사과 검출 수가 ~24까지 단조 증가.",
            (1, "apple_1 (GT 48): anchor 8→24 변경 시 예측 22 → 28."),
            (1, "배(상대적으로 작음)는 변화 없음."),
            "SAM이 스케일에 민감한 이유:",
            (1, "SAM은 객체 INSIDE 프롬프트에서 segment하도록 학습됨. 작은 anchor도 사과 내부에 있지만, 마스크 디코더의 best 가설이 anchor 스케일을 따라감."),
            (1, "연구 질문: 스케일 적응형 anchor 앙상블? 다중 스케일 투표?"),
        ],
    )

    add_content_slide(
        prs,
        "인사이트 6 — 정렬을 깨뜨리는 파인튜닝보다 프롬프트 엔지니어링이 강할 수 있음",
        [
            "과수원 이미지에서 upstream MLP의 프롬프트별 MAE:",
            (1, "`pears` → 35.0 · `pear in paper bag` → 6.0"),
            (1, "`apples` → 42.5 · `apple in paper bag` → 22.5"),
            (1, "`protective fruit bag` → 40.0"),
            "배에서 세 단어의 프롬프트 변경이 MAE 약 30을 줄임. 재학습 필요 없음.",
            "메커니즘: CLIP 텍스트 인코더가 조합적 의미(과일 종류 × 포장 맥락)를 포착.",
            (1, "\"apple in paper bag\" — 사과 클러스터 + 포장 맥락을 동시 활성화."),
            (1, "\"apples\" 단독은 사과 클러스터만 활성화 — 덜 구체적 → 넓은 장면에서 과잉 검출."),
            "함의: 클래스 무관 카운팅 파이프라인에서는 파인튜닝이 아니라 프롬프트 서치를 FIRST 개입으로.",
            (1, "파인튜닝은 CLIP 정렬 손실 위험. 프롬프트 서치는 공짜이고 가역적."),
        ],
    )

    # =============================================================
    # Part 4: 연구 방향
    # =============================================================
    add_section_slide(prs, 4, "열린 질문 · 향후 연구")

    add_two_col_slide(
        prs,
        "이번 이터레이션에서 답하지 못한 질문",
        "데이터 의존적",
        [
            "봉지 과일 실제 이미지에서 우리 기본 파이프라인의 실제 성능은?",
            (1, "우리가 받은 captures는 모두 흰 픽셀(255)로 비어 있었음."),
            "도메인 라벨 없음. 50장 정도 수동 카운팅하면 정상적 MAE 곡선 가능.",
            "남은 사과 undercount (GT ~47, pred ~30)는 SAM 스케일 문제인가, CLIP 의미 문제인가?",
            (1, "다중 스케일 anchor 앙상블 ablation 필요."),
        ],
        "이론 의존적",
        [
            "단일 보조 손실로 파인튜닝 중 CLIP 정렬을 유지할 수 있는가?",
            (1, "PseCo의 `cls_loss2`는 후보 — 별도 배치에서 CLIP-검출 region에 대한 일관성 항."),
            (1, "더 단순한 text-prompt-discrimination 손실로 대체 가능할까?"),
            "스케일 적응형 proposal 생성 — 현재는 도메인별 수동 anchor 튜닝.",
            "\"파인튜닝 후에도 CLIP 정렬 유지\"를 측정하는 원리적 지표는 무엇인가?",
        ],
    )

    add_content_slide(
        prs,
        "권장 연구 방향",
        [
            "Direction A — 정렬 보존 파인튜닝.",
            (1, "PseCo의 `cls_loss2` (CLIP-검출 region에 대한 별도 배치 일관성 손실)을 우리 트레이너로 이식."),
            (1, "정렬을 잃은 Plan 3 트레이너와 비교 평가."),
            (1, "목표: 과수원 데이터 파인튜닝 중에도 prompt-discrimination을 임계 이상으로 유지."),
            "Direction B — 스케일 적응형 카운팅.",
            (1, "다중 anchor 앙상블: anchor ∈ {4, 8, 16, 24, 32} 병렬 실행 후 점수 통합 + 공동 NMS."),
            (1, "혹은 이미지 통계로부터 anchor 크기를 예측하는 작은 모델 학습."),
            "Direction C — Baseline 인식 카운팅 지표.",
            (1, "MAE - predict-mean, MAE - predict-zero 두 값을 함께 리포트."),
            (1, "왜곡된 카운트 분포에서 \"실제 학습\"과 \"분포 암기\"를 구분해 줌."),
            "Direction D — 프롬프트 최적화를 일급 단계로.",
            (1, "파인튜닝 고려 전에 체계적 프롬프트 서치를 수행."),
            (1, "이미지 임베딩에 조건화된 학습형 프롬프트 생성기로 하이브리드화 가능."),
        ],
    )

    add_content_slide(
        prs,
        "향후 진행 계획",
        [
            "Plan 5 (즉시 재개 조건부) — 봉지 과일 데이터 확보 + 기초 검증",
            (1, "작업: 배포 카메라로 봉지 과일 10~30장 재촬영 (현재 샘플은 흰 픽셀로 결함)."),
            (1, "작업: 수동 라벨링 — 이미지당 총 개수 + 핵심 샘플은 점 좌표까지."),
            (1, "검증: upstream MLP 기본 설정 + 기본/튜닝 프롬프트로 MAE 측정. Plan 4 범위 내 해결 가능 여부 판정."),
            "Plan 6 (조건부) — 정렬 보존 파인튜닝(C3)",
            (1, "진입 조건: Plan 5 검증에서 upstream MAE가 불충분 (예: MAE > 10)."),
            (1, "작업: upstream `3_extract_proposals.py`의 pseudo-box 생성 이식."),
            (1, "작업: `cls_loss2` (CLIP 일관성 손실) 로 듀얼 loss 학습 루프 구현."),
            (1, "작업: 과수원 라벨로 점진적 파인튜닝 + prompt-discrimination 모니터링."),
            "Plan 7 (선택) — 원래 Plan 4 스펙의 미사용 아이템",
            (1, "Classifier (ResNet18) 학습 — PseCo만으로 부족한 경우 verify 단계 추가."),
            (1, "Super-Resolution 통합 — small crop에서의 classification 개선 여지."),
            (1, "과수원 전용 UI / 결과 시각화 대시보드 (비즈니스 요구 시)."),
        ],
        footer="Plan 5 시작 선결조건: 실제 봉지 과일 이미지 + 라벨. 그 이전까지는 Plan 4 기본이 프로덕션.",
    )

    add_table_slide(
        prs,
        "예상 결과 (Target Metrics)",
        headers=["지표", "현재 (Plan 4 기본)", "Plan 5 목표", "Plan 6 목표 (파인튜닝 후)"],
        rows=[
            ("과수원 phone 이미지 MAE (4장)", "11.75", "≤ 10", "≤ 5"),
            ("사과 undercount (GT 46-48)", "pred 28-30 (-35%)", "pred 35+ (-25%)", "pred 42+ (-10%)"),
            ("배 MAE", "3-8", "< 5 (유지)", "< 3"),
            ("봉지 과일 MAE", "측정 불가(데이터 없음)", "≤ 15 (라벨 확보 후)", "≤ 8"),
            ("CLIP 프롬프트 감도", "upstream 수준 유지", "유지", "정렬 보존 학습으로 유지"),
            ("val/mae on FSC-147", "파인튜닝 미적용", "—", "≤ 20 (upstream 수준)"),
            ("추론 속도 (RTX 5070)", "이미지당 ~2 초", "유지", "유지"),
        ],
        col_widths=[3.3, 3.0, 3.3, 3.3],
        footer="모든 숫자는 현 데이터 분포 기반 추정치. 실제 결과는 확보한 봉지 이미지의 품질·다양성에 따라 변동.",
    )

    add_content_slide(
        prs,
        "진입 조건 및 타임라인 (개략)",
        [
            "Plan 5 진입 조건: 정상적으로 내용물이 있는 봉지 과일 이미지 최소 10장 + 총 개수 라벨.",
            (1, "조건 충족 시 예상 소요: ~1일 (데이터 업로드 0.5시간 + 평가 스크립트 실행 0.5시간 + 결과 분석)."),
            "Plan 6 진입 조건: Plan 5 결과에서 upstream + 튜닝만으로 MAE 목표 미달.",
            (1, "`cls_loss2` 이식 + pseudo-box 생성 + 학습 루프 재작성: ~3~5일 예상."),
            (1, "실제 파인튜닝 + 검증: 추가 1~2일 (RTX 5070 기준)."),
            "Plan 7 진입 조건: PseCo 단독으로 false positive(봉지 아닌 것 카운팅) 가 문제일 때만.",
            (1, "ResNet18 분류기 학습: 라벨된 \"봉지 crop\" 500장 이상 필요."),
            "전체 일정은 라벨링 속도에 가장 크게 의존. 수작업 라벨링 페이스에 따라 전체 기간 결정.",
        ],
    )

    add_content_slide(
        prs,
        "리스크 및 감지 지표 (학습 중 체크포인트)",
        [
            "리스크 R1 — Plan 5에서 실제 봉지 이미지 품질이 낮음 (흐림, 원거리, 조명 불량).",
            (1, "감지: 진단 툴(`counting diagnose`)의 blur/exposure 분포로 선제 확인."),
            (1, "대응: 촬영 가이드라인 재정립 혹은 이미지 Deblur/SR 전처리 활성화."),
            "리스크 R2 — Plan 6 파인튜닝에서 CLIP 정렬을 다시 잃음 (Plan 3 재현).",
            (1, "감지: 매 N iter마다 prompt-discrimination 점수 (fixed triplet set) 기록."),
            (1, "대응: `cls_loss2` 가중치 상향 / early stopping 기준 강화."),
            "리스크 R3 — 사과 undercount가 SAM 스케일이 아닌 더 깊은 문제 (가림, 색상 유사성).",
            (1, "감지: PointDecoder의 점 수 + IoU(pred_box, visible_apples) 시각화."),
            (1, "대응: 데이터 증강 (rotation, partial occlusion), 혹은 upstream 대신 다른 proposal 생성 방식."),
            "리스크 R4 — GPU 경합 (사용자의 다른 학습과 동시 실행).",
            (1, "감지: `nvidia-smi --query-gpu=memory.used`."),
            (1, "대응: 학습을 순차 실행, 배치 사이즈 동적 조정."),
        ],
    )

    add_content_slide(
        prs,
        "엔지니어링 측면 교훈 (연구는 아니지만 남길 가치 있음)",
        [
            "캐시의 무효화 인식 설계는 시간을 크게 절약함. SAM 피처는 `(ckpt 해시, 이미지 크기, dtype)` 해시로 저장.",
            "모든 ML 파이프라인에 'predict-0' 참조 런을 함께 돌려야 함. 안 그러면 데이터셋 평균 모사에 속아 축하할 수 있음.",
            "모든 체크포인트 안에 설정 스냅샷을 같이 저장할 것. \"이 ckpt가 더 이상 로드 안 됨\" 순간이 몇 번 있었음.",
            "`slow` 마커 + `skipif(가중치 없음)` 붙은 통합 테스트는 없는 것보단 낫지만 주기적으로 실제로 돌려야 함.",
            "Upstream 레포를 submodule로 받을 땐 \"demo.ipynb가 실제로 깨끗이 import되는가\"를 20분 안에 점검하고 API surface를 믿을지 결정.",
            (1, "우리 `PseCoStage.prepare()`는 세 Plan 동안 존재하지 않는 모듈을 import하고 있었음."),
        ],
    )

    add_content_slide(
        prs,
        "참고문헌",
        [
            "PseCo — Huang et al., \"Point, Segment, and Count: A Generalized Framework for Object Counting\", ICLR 2024. github.com/Hzzone/PseCo",
            "SAM — Kirillov et al., \"Segment Anything\", ICCV 2023. segment-anything.com",
            "CLIP — Radford et al., \"Learning Transferable Visual Models From Natural Language Supervision\", ICML 2021.",
            "OpenCLIP — Ilharco et al. github.com/mlfoundations/open_clip",
            "FSC-147 — Ranjan et al., \"Learning To Count Everything\", CVPR 2021. github.com/cvlab-stonybrook/LearningToCountEverything",
            "본 프로젝트 — github.com/YBNML/fruit-counting",
            (1, "스펙: docs/superpowers/specs/"),
            (1, "Plan별 구현 기록: docs/superpowers/plans/"),
        ],
    )

    prs.save(str(OUT))
    print(f"[done] wrote {OUT}")


if __name__ == "__main__":
    main()
